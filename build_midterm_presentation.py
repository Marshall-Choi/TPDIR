"""
EEE429 Mid-term PPTX builder — based on EEE426_FINAL.ipynb actual outputs.
Run: .venv/bin/python build_midterm_presentation.py
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── colour palette ──────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1E, 0x1E, 0x2E)   # near-black background
C_ACCENT = RGBColor(0x74, 0xC7, 0xEC)   # cornflower blue
C_GREEN  = RGBColor(0xA6, 0xE3, 0xA1)
C_YELLOW = RGBColor(0xF9, 0xE2, 0xAF)
C_RED    = RGBColor(0xF3, 0x8B, 0xA8)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xCD, 0xD6, 0xF4)
C_GRAY   = RGBColor(0x58, 0x5B, 0x70)
C_SURFACE= RGBColor(0x31, 0x32, 0x44)

W = Inches(13.333)
H = Inches(7.5)

OUT = Path(__file__).resolve().parent / "EEE429_Midterm_Presentation.pptx"

# ── helpers ──────────────────────────────────────────────────────────────────

def _bg(slide, color: RGBColor = C_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, text: str, left, top, width, height, *,
         font_size=20, bold=False, color=C_WHITE, bg=None,
         align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    return txBox


def _rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _title(slide, text: str, subtitle: str = ""):
    _rect(slide, 0, 0, W, Inches(1.35), C_SURFACE)
    _box(slide, text, Inches(0.5), Inches(0.18), Inches(12), Inches(0.9),
         font_size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        _box(slide, subtitle, Inches(0.5), Inches(1.02), Inches(12), Inches(0.38),
             font_size=16, color=C_LIGHT, align=PP_ALIGN.LEFT)


def _note(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def _table(slide, rows, left, top, width, height,
           header_bg=C_ACCENT, row_bg=C_SURFACE, alt_bg=C_DARK,
           header_color=C_DARK, cell_color=C_WHITE,
           col_widths: list[float] | None = None):
    """rows[0] = header row, rows[1:] = data rows. Each cell is a string."""
    nrows = len(rows)
    ncols = len(rows[0])
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for ci in range(ncols):
        if col_widths:
            tbl.columns[ci].width = int(width * col_widths[ci])
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(14 if ri == 0 else 13)
            run.font.bold = (ri == 0)
            run.font.color.rgb = header_color if ri == 0 else cell_color
            fill = cell.fill
            fill.solid()
            if ri == 0:
                fill.fore_color.rgb = header_bg
            else:
                fill.fore_color.rgb = row_bg if ri % 2 == 1 else alt_bg
    return tbl


# ── slides ───────────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(s)
    _rect(s, 0, Inches(2.5), W, Inches(2.6), C_SURFACE)
    _box(s, "EEE429 Final Project", Inches(1), Inches(2.6),
         Inches(11), Inches(0.7), font_size=44, bold=True,
         color=C_ACCENT, align=PP_ALIGN.CENTER)
    _box(s, "Mid-Term Progress Presentation",
         Inches(1), Inches(3.3), Inches(11), Inches(0.6),
         font_size=26, color=C_WHITE, align=PP_ALIGN.CENTER)
    _box(s, "MNIST CNN — Quantization-Aware Training & HW Artifact Generation",
         Inches(1), Inches(3.9), Inches(11), Inches(0.5),
         font_size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _box(s, "Based on: EEE426_FINAL.ipynb",
         Inches(1), Inches(5.2), Inches(11), Inches(0.4),
         font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
    _note(s,
        "Good morning. This is our mid-term progress report for EEE429. "
        "We will walk through exactly what we have completed so far — "
        "the quantized CNN model, training results, and all hardware artifact exports — "
        "as implemented in EEE426_FINAL.ipynb.")


def slide_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "What We Have Completed",
           "EEE426_FINAL.ipynb — full software pipeline end-to-end")

    items = [
        ("①", "Quantization Scheme",
         "ap_fixed<16,7> emulated in PyTorch with clamp + round"),
        ("②", "Quantized CNN Model",
         "QUANTIZED_MNISTCNN — full 3-Conv + Pool + FC, QAT throughout"),
        ("③", "Training & Evaluation",
         "5 epochs · Adam · batch 256 · test acc peaked at 98.92%"),
        ("④", "Best Checkpoint",
         "QUANTIZED_MNISTCNN-2026-05-19-00-18-25.pth · full-set acc 99.36%"),
        ("⑤", "weights.h Export",
         "All 8 tensors snapped to fixed-point grid, timestamped C header"),
        ("⑥", "images.h Export",
         "280 images, shape [280][1][28][28], float C header"),
        ("⑦", "MNIST .npy Generation",
         "70 k images + labels digit-sorted, ready for co-sim"),
    ]
    top = Inches(1.55)
    step = Inches(0.77)
    for num, title, desc in items:
        _rect(s, Inches(0.35), top, Inches(0.5), Inches(0.52), C_ACCENT)
        _box(s, num, Inches(0.35), top, Inches(0.5), Inches(0.52),
             font_size=15, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        _box(s, title, Inches(1.0), top, Inches(3.5), Inches(0.52),
             font_size=15, bold=True, color=C_YELLOW)
        _box(s, desc, Inches(4.6), top, Inches(8.4), Inches(0.52),
             font_size=14, color=C_LIGHT)
        top += step

    _note(s,
        "This slide summarises all completed deliverables. "
        "Every item listed has been executed and outputs verified inside the notebook. "
        "I will go through each in detail on the following slides.")


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "CNN Architecture", "Spec-compliant topology — unchanged from project brief")

    # diagram boxes (left to right)
    layers = [
        ("Input\n1×28×28",       C_GRAY),
        ("Conv1\n1→16, k=3\n→16×26×26",   C_ACCENT),
        ("Conv2\n16→32, k=3\n→32×24×24",  C_ACCENT),
        ("Conv3\n32→32, k=3\n→32×22×22",  C_ACCENT),
        ("MaxPool\n2×2\n→32×11×11",        C_GREEN),
        ("Flatten\n3 872",                  C_SURFACE),
        ("FC\n3872→10",                     C_YELLOW),
        ("Output\nLabel",                   C_GRAY),
    ]
    bw = Inches(1.5)
    bh = Inches(1.3)
    gap = Inches(0.18)
    total = len(layers) * bw + (len(layers)-1)*gap
    start_x = (W - total) / 2
    y = Inches(2.1)
    for i, (lbl, col) in enumerate(layers):
        x = start_x + i*(bw+gap)
        _rect(s, x, y, bw, bh, col)
        _box(s, lbl, x, y, bw, bh, font_size=11, bold=False,
             color=C_DARK if col in (C_ACCENT, C_GREEN, C_YELLOW) else C_WHITE,
             align=PP_ALIGN.CENTER)
        if i < len(layers)-1:
            _box(s, "→", x+bw, y+Inches(0.45), gap+Inches(0.05), Inches(0.45),
                 font_size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # quantization annotations
    annots = [
        (start_x + 0*(bw+gap), "Q"),
        (start_x + 1*(bw+gap), "Q"),
        (start_x + 2*(bw+gap), "Q"),
        (start_x + 3*(bw+gap), "Q"),
        (start_x + 4*(bw+gap), "Q"),
    ]
    for ax, label in annots:
        _rect(s, ax + Inches(0.55), y - Inches(0.38), Inches(0.4), Inches(0.3), C_RED)
        _box(s, label, ax + Inches(0.55), y - Inches(0.38), Inches(0.4), Inches(0.3),
             font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _box(s, "Q = FixedPointQuantize inserted before ReLU at each layer output + after pool",
         Inches(0.5), Inches(3.65), Inches(12), Inches(0.35),
         font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

    # param table
    rows = [
        ["Layer", "Class", "In → Out channels", "Kernel / Stride", "Output shape", "Bias"],
        ["Conv1", "QuantizedConv2d", "1 → 16", "3×3 / 1", "16 × 26 × 26", "✓"],
        ["Conv2", "QuantizedConv2d", "16 → 32", "3×3 / 1", "32 × 24 × 24", "✓"],
        ["Conv3", "QuantizedConv2d", "32 → 32", "3×3 / 1", "32 × 22 × 22", "✓"],
        ["MaxPool", "nn.MaxPool2d", "—", "2×2", "32 × 11 × 11", "—"],
        ["FC", "QuantizedLinear", "3 872 → 10", "—", "10", "✓"],
    ]
    _table(s, rows, Inches(0.3), Inches(4.15), Inches(12.73), Inches(3.0),
           col_widths=[0.09, 0.18, 0.17, 0.16, 0.17, 0.07])

    _note(s,
        "The architecture is fixed by the project spec. "
        "All three convolutional layers use valid-padding three-by-three kernels, stride one. "
        "The spatial dimension steps from twenty-eight to twenty-six to twenty-four to twenty-two. "
        "After the two-by-two max pool we get eleven by eleven by thirty-two, which flattens to 3872. "
        "The FC layer maps that to ten class scores. "
        "Red Q badges show where a fixed-point quantizer is inserted in the forward pass.")


def slide_fixedpoint(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "Fixed-Point Scheme: ap_fixed⟨16, 7⟩",
           "Emulated in Python; identical semantics to Vivado HLS ap_fixed")

    _rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(2.4), C_SURFACE)
    specs = [
        ("Word width   W", "16 bits"),
        ("Integer bits  I", "7  bits  (includes sign)"),
        ("Fraction bits F", "9  bits  (= W − I)"),
        ("Scale factor", "2⁹ = 512"),
        ("Min value", "−64"),
        ("Max value", "64 − 2⁻⁹  ≈  63.998"),
        ("Quantum / LSB", "2⁻⁹ ≈ 0.00195"),
    ]
    y = Inches(1.6)
    for k, v in specs:
        _box(s, k, Inches(0.6), y, Inches(2.6), Inches(0.3),
             font_size=13, color=C_LIGHT)
        _box(s, v, Inches(3.3), y, Inches(2.7), Inches(0.3),
             font_size=13, bold=True, color=C_YELLOW)
        y += Inches(0.3)

    _box(s, "16-bit word layout",
         Inches(6.5), Inches(1.55), Inches(6.5), Inches(0.35),
         font_size=13, bold=True, color=C_LIGHT)
    bits_y = Inches(1.95)
    bw_bit = Inches(0.4)
    for i in range(16):
        col = C_RED if i == 0 else (C_ACCENT if i < 7 else C_GREEN)
        _rect(s, Inches(6.5) + i*bw_bit, bits_y, bw_bit - Inches(0.02), Inches(0.5), col)
        label = "S" if i == 0 else (str(6-i) if i < 7 else f"-{i-6}")
        _box(s, label, Inches(6.5) + i*bw_bit, bits_y, bw_bit, Inches(0.5),
             font_size=9, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)
    _box(s, "← Integer (7 bits, incl. sign) →",
         Inches(6.5), bits_y+Inches(0.53), Inches(2.8), Inches(0.3),
         font_size=10, color=C_RED, align=PP_ALIGN.CENTER)
    _box(s, "← Fraction (9 bits) →",
         Inches(9.3), bits_y+Inches(0.53), Inches(3.6), Inches(0.3),
         font_size=10, color=C_GREEN, align=PP_ALIGN.CENTER)

    code = (
        "def fixed_point_round(x):\n"
        "    x = torch.clamp(x, -64, 64 - 2⁻⁹)\n"
        "    x = torch.round(x * 512) / 512\n"
        "    return x\n\n"
        "# STE backward\n"
        "def backward(ctx, grad):\n"
        "    return grad  # pass-through"
    )
    _rect(s, Inches(0.4), Inches(4.1), Inches(12.5), Inches(2.4), C_SURFACE)
    _box(s, code, Inches(0.55), Inches(4.2), Inches(12.2), Inches(2.2),
         font_size=13, color=C_GREEN)

    _note(s,
        "We use sixteen-bit signed fixed-point with seven integer bits and nine fraction bits. "
        "The representable range is minus sixty-four to roughly plus sixty-four. "
        "In Python the emulator clamps the tensor to that range and rounds to the nearest 1/512. "
        "Importantly the backward pass uses a straight-through estimator — the gradient passes through the "
        "quantizer unchanged so training stays stable. "
        "The exported C header declares using fixed_t equal to ap_fixed sixteen seven AP_RND AP_SAT "
        "which is bitwise identical to this emulation.")


def slide_qat_forward(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "QAT Forward Graph",
           "Exact quantization insertion order from QUANTIZED_MNISTCNN.forward()")

    steps = [
        ("INPUT x",           C_GRAY,    "Normalize(0.5, 0.5) → range [−1, +1]"),
        ("Q(x)",              C_RED,     "FixedPointQuantize  — input quantization"),
        ("QuantizedConv2d 1", C_ACCENT,  "q_w = Q(self.weight),  q_b = Q(self.bias)"),
        ("Q(conv1 out)",      C_RED,     "FixedPointQuantize  — before ReLU"),
        ("ReLU",              C_GREEN,   "in-place clamp ≥ 0"),
        ("QuantizedConv2d 2", C_ACCENT,  "same pattern"),
        ("Q(conv2 out)",      C_RED,     "FixedPointQuantize"),
        ("ReLU",              C_GREEN,   ""),
        ("QuantizedConv2d 3", C_ACCENT,  ""),
        ("Q(conv3 out)",      C_RED,     "FixedPointQuantize"),
        ("ReLU",              C_GREEN,   ""),
        ("MaxPool2d(2)",       C_SURFACE, "kernel 2×2"),
        ("Q(pool out)",       C_RED,     "FixedPointQuantize  — after pool"),
        ("Flatten → FC",      C_YELLOW,  "QuantizedLinear(3872→10)"),
        ("Logits",            C_GRAY,    "CrossEntropyLoss during training"),
    ]
    cols = 3
    bw = Inches(4.1)
    bh = Inches(0.38)
    gap_x = Inches(0.1)
    gap_y = Inches(0.04)
    start_x = Inches(0.25)
    start_y = Inches(1.55)
    for i, (label, col, desc) in enumerate(steps):
        row = i // cols
        c   = i  % cols
        x = start_x + c * (bw + gap_x)
        y = start_y + row * (bh + gap_y)
        _rect(s, x, y, bw, bh, col)
        full = f"{label}" + (f"  —  {desc}" if desc else "")
        _box(s, full, x + Inches(0.08), y, bw - Inches(0.08), bh,
             font_size=11,
             color=C_DARK if col in (C_ACCENT, C_GREEN, C_YELLOW) else C_WHITE)

    _note(s,
        "This is the exact forward pass order in the notebook. "
        "The input is first snapped to the fixed-point grid. "
        "Then for each convolution we quantize weights and biases inside QuantizedConv2d, "
        "quantize the convolution output, and then apply ReLU. "
        "After the third convolution and ReLU the pooled output is quantized one more time. "
        "Finally the flattened vector goes through QuantizedLinear which also quantizes its weights. "
        "The HLS kernel must mirror this order exactly.")


def slide_training_setup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "Training Setup", "Hyperparameters and environment — from notebook cells")

    rows = [
        ["Parameter",            "Value"],
        ["Framework",            "PyTorch + torchvision  (Colab, T4 GPU)"],
        ["Model class",          "QUANTIZED_MNISTCNN"],
        ["Optimizer",            "Adam"],
        ["Learning rate",        "1 × 10⁻³"],
        ["Batch size",           "256"],
        ["Epochs",               "5"],
        ["Loss function",        "CrossEntropyLoss"],
        ["Input preprocessing",  "ToTensor  +  Normalize(mean=0.5, std=0.5)  →  [−1, +1]"],
        ["Data split",           "train: 60 000   test: 10 000   total eval: 70 000"],
        ["Checkpoint saves",     "Best test accuracy (each epoch saved if improved)"],
    ]
    _table(s, rows, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.6),
           col_widths=[0.28, 0.72])

    _note(s,
        "Training ran on Google Colab with a T4 GPU. "
        "We used Adam with one-thousandth learning rate, batch size two fifty-six, and five epochs. "
        "Best checkpoint is saved to milestones. "
        "Input preprocessing normalises pixel values from zero-one to minus-one to plus-one. "
        "This same transform must be applied on the board driver at inference time.")


def slide_training_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "Training Results — Epoch Log",
           "QUANTIZED_MNISTCNN · 5 epochs · Adam lr=1e-3 · batch=256")

    rows_train = [
        ["Epoch", "Train Acc", "Train Loss", "Test Acc", "Test Loss", "Saved?"],
        ["0", "91.69 %", "0.29000", "97.97 %", "0.06443", "✓  *"],
        ["1", "98.16 %", "0.06245", "98.44 %", "0.04900", "✓  *"],
        ["2", "98.62 %", "0.04626", "98.69 %", "0.04052", "✓  *"],
        ["3", "98.87 %", "0.03645", "98.64 %", "0.03848", "—"],
        ["4", "99.05 %", "0.03035", "98.92 %", "0.03147", "✓  *  ← best"],
    ]
    _table(s, rows_train, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.75),
           col_widths=[0.08, 0.15, 0.15, 0.15, 0.15, 0.20])

    _box(s, "Full-dataset evaluation (train + test = 70 000 samples)",
         Inches(0.5), Inches(4.5), Inches(12), Inches(0.35),
         font_size=15, bold=True, color=C_YELLOW)

    rows_full = [
        ["Checkpoint filename", "Full-set Accuracy"],
        ["QUANTIZED_MNISTCNN-2026-05-19-00-18-25.pth  (latest)", "99.36 %"],
        ["QUANTIZED_MNISTCNN-2026-05-17-23-48-02.pth  (previous)", "99.10 %"],
    ]
    _table(s, rows_full, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.55),
           col_widths=[0.72, 0.28])

    _note(s,
        "Looking at the epoch log: training accuracy climbs from ninety-one percent to ninety-nine percent over five epochs. "
        "Test accuracy peaked at ninety-eight point ninety-two percent on epoch four, which triggered a checkpoint save. "
        "The full-dataset number of ninety-nine point thirty-six percent is higher because the model was also evaluated "
        "on its own training data. The key number for the course is the test accuracy of ninety-eight point ninety-two percent "
        "which is well above the ninety-three percent threshold.")


def slide_weight_embed(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "WeightEmbed — weights.h Export",
           "Checkpoint class → ap_fixed C header, auto-timestamped")

    # left: pipeline
    steps = [
        ("Load .pth",  "torch.load(path, map_location='cpu')"),
        ("Iterate state_dict", "keys: conv1.weight / .bias … fc.weight / .bias"),
        ("fixed_point_round_tensor", "clamp → round to 1/512 grid"),
        ("_nested_format", "Recursively build C brace-initialiser string"),
        ("Write .h", "using fixed_t = ap_fixed<16,7,AP_RND,AP_SAT>"),
    ]
    bh = Inches(0.58)
    y = Inches(1.62)
    for title, desc in steps:
        _rect(s, Inches(0.35), y, Inches(0.22), bh, C_ACCENT)
        _box(s, title, Inches(0.65), y, Inches(3.2), bh,
             font_size=13, bold=True, color=C_ACCENT)
        _box(s, desc, Inches(3.9), y, Inches(3.4), bh,
             font_size=12, color=C_LIGHT)
        if y < Inches(4.5):
            _rect(s, Inches(0.42), y+bh, Inches(0.08), Inches(0.1), C_ACCENT)
        y += bh + Inches(0.1)

    # right: tensor table
    rows = [
        ["Tensor key", "Shape", "Total params"],
        ["conv1.weight", "(16, 1, 3, 3)", "144"],
        ["conv1.bias",   "(16,)",          "16"],
        ["conv2.weight", "(32, 16, 3, 3)", "4 608"],
        ["conv2.bias",   "(32,)",          "32"],
        ["conv3.weight", "(32, 32, 3, 3)", "9 216"],
        ["conv3.bias",   "(32,)",          "32"],
        ["fc.weight",    "(10, 3872)",     "38 720"],
        ["fc.bias",      "(10,)",          "10"],
        ["Total",        "—",             "52 778"],
    ]
    _table(s, rows, Inches(7.55), Inches(1.62), Inches(5.5), Inches(5.2),
           col_widths=[0.4, 0.35, 0.25])

    _box(s,
         "Output: weights_YYYY-MM-DD-HH-MM-SS.h  ·  "
         "Latest: weights_2026-05-19-00-21-15.h  (7 983 lines)",
         Inches(0.35), Inches(6.55), Inches(12.6), Inches(0.4),
         font_size=12, color=C_GRAY)

    _note(s,
        "The Checkpoint class loads the PyTorch state dictionary and iterates over every tensor. "
        "Each tensor is snapped to the sixteen-seven fixed-point grid using the same clamp-round function as training. "
        "It then generates a nested C brace initialiser and writes one self-contained header file. "
        "The file declares a using typedef matching the HLS ap_fixed type so the hardware C code can include it directly. "
        "Shape defines such as CONV1_WEIGHT_DIM0 and so on are emitted for each tensor so the HLS kernel can "
        "reference them symbolically. The latest output file is about eight thousand lines long.")


def slide_image_embed(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "ImageEmbed — images.h & MNIST NumPy Export",
           "Test vectors for HLS co-simulation and hardware bring-up")

    # images.h section
    _rect(s, Inches(0.35), Inches(1.55), Inches(6.1), Inches(2.5), C_SURFACE)
    _box(s, "images.h", Inches(0.5), Inches(1.6),
         Inches(5), Inches(0.4), font_size=18, bold=True, color=C_ACCENT)
    details_h = [
        "280 sample images, 28 per digit (digits 0–9)",
        "Shape: [280][1][28][28]  — NOT flattened",
        "dtype: float  (same Normalize(0.5, 0.5) as training)",
        "Format: static const float IMAGES[280][1][28][28] = {…}",
        "#define IMAGE_NUM 280",
    ]
    y = Inches(2.1)
    for d in details_h:
        _box(s, "• " + d, Inches(0.6), y, Inches(5.7), Inches(0.32),
             font_size=12, color=C_LIGHT)
        y += Inches(0.33)

    # npy section
    _rect(s, Inches(6.7), Inches(1.55), Inches(6.3), Inches(2.5), C_SURFACE)
    _box(s, "MNIST_DATASET_*.npy", Inches(6.85), Inches(1.6),
         Inches(6), Inches(0.4), font_size=18, bold=True, color=C_GREEN)
    details_npy = [
        "Full 70 000 images  +  labels",
        "Sorted by digit: 0×7000, 1×7000 … 9×7000",
        "IMAGE shape: (70000, 1, 28, 28)  — preserves CHW",
        "LABEL shape: (70000,)  int64",
        "Used for RTL/HLS co-simulation ground truth",
    ]
    y = Inches(2.1)
    for d in details_npy:
        _box(s, "• " + d, Inches(6.85), y, Inches(6.0), Inches(0.32),
             font_size=12, color=C_LIGHT)
        y += Inches(0.33)

    # code snippet
    _box(s, "Export calls (from notebook):",
         Inches(0.35), Inches(4.2), Inches(12), Inches(0.3),
         font_size=13, bold=True, color=C_YELLOW)
    code = (
        'mnist.export_h(280, "float", lambda x: f"{repr(x)}f", "./images.h")\n'
        'mnist.export_np("./MNIST_DATASET_IMAGE.npy", mode="image")\n'
        'mnist.export_np("./MNIST_DATASET_LABEL.npy", mode="label")'
    )
    _rect(s, Inches(0.35), Inches(4.55), Inches(12.6), Inches(1.0), C_SURFACE)
    _box(s, code, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.9),
         font_size=13, color=C_GREEN)

    _box(s,
         "Note: MNIST data is digit-bucketed, then interleaved 28 samples per digit for images.h",
         Inches(0.35), Inches(5.7), Inches(12.6), Inches(0.3),
         font_size=11, color=C_GRAY)

    _note(s,
        "The image export section produces two types of output. "
        "images.h contains two hundred eighty hand-picked samples — twenty-eight per digit — "
        "in a four-dimensional C array. Shape is preserved as batch by channel by height by width. "
        "This allows the hardware team to compile a quick functional test without extra infrastructure. "
        "The NumPy exports cover all seventy thousand MNIST images sorted by digit class, "
        "which is convenient for co-simulation where you want to sweep a full digit at a time.")


def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "Complete Software Pipeline",
           "EEE426_FINAL.ipynb — end-to-end flow")

    steps = [
        ("DATASET", "MNIST 70k\nToTensor + Normalize(0.5,0.5)",  C_GRAY),
        ("QAT MODEL", "QUANTIZED_MNISTCNN\nap_fixed<16,7> throughout", C_ACCENT),
        ("TRAINING", "Adam, lr=1e-3\nbatch 256, 5 epochs",          C_ACCENT),
        ("BEST CKPT", "QUANTIZED_MNISTCNN-\n2026-05-19-00-18-25.pth", C_GREEN),
        ("weights.h", "8 tensors, C header\nap_fixed<16,7,AP_RND,AP_SAT>", C_YELLOW),
        ("images.h", "280 images\n[280][1][28][28] float",           C_YELLOW),
        ("MNIST .npy", "IMAGE (70k,1,28,28)\nLABEL (70k,) int64",   C_YELLOW),
    ]

    bw = Inches(1.7)
    bh = Inches(1.4)
    gx = Inches(0.22)
    total = len(steps)*bw + (len(steps)-1)*gx
    sx = (W - total)/2
    y  = Inches(2.0)

    for i, (title, desc, col) in enumerate(steps):
        x = sx + i*(bw+gx)
        _rect(s, x, y, bw, bh, col)
        _box(s, title, x, y+Inches(0.1), bw, Inches(0.42),
             font_size=13, bold=True,
             color=C_DARK if col != C_GRAY else C_WHITE,
             align=PP_ALIGN.CENTER)
        _box(s, desc, x, y+Inches(0.52), bw, Inches(0.82),
             font_size=10,
             color=C_DARK if col != C_GRAY else C_LIGHT,
             align=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            ax = x+bw+gx*0.1
            _box(s, "→", ax, y+Inches(0.5), gx+Inches(0.1), Inches(0.5),
                 font_size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # accuracy annotation
    _rect(s, Inches(3.6), Inches(3.65), Inches(2.5), Inches(0.55), C_SURFACE)
    _box(s, "Test: 98.92 %\nFull-set: 99.36 %",
         Inches(3.65), Inches(3.68), Inches(2.4), Inches(0.5),
         font_size=11, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    # divider
    _rect(s, Inches(7.0), Inches(1.65), Inches(0.04), Inches(2.2), C_GRAY)
    _box(s, "SW output for\nhardware team →",
         Inches(7.15), Inches(1.95), Inches(1.5), Inches(0.7),
         font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

    _box(s, "All outputs verified and committed. Ready to hand off to hardware team.",
         Inches(0.35), Inches(5.55), Inches(12.6), Inches(0.4),
         font_size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    _note(s,
        "This diagram shows the full software pipeline as it stands today. "
        "Starting from MNIST data we train QUANTIZED_MNISTCNN with QAT. "
        "The best checkpoint at ninety-eight point ninety-two percent test accuracy is saved. "
        "From that checkpoint we export a weights.h file containing all eight tensors "
        "snapped to the fixed-point grid, plus images.h and NumPy archives for co-simulation. "
        "Everything to the right of the dotted line is an artifact consumed by the hardware team.")


def slide_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _title(s, "Remaining Work & Interface Contract",
           "What still needs to be agreed / implemented")

    items_hw = [
        "HLS kernel implementing exact forward graph (conv→Q→relu order)",
        "DMA / Overlay integration on Pynq-Z2",
        "Bitstream (.bit) + hardware handoff (.hwh)",
        "C-sim / Co-sim reports using exported weights.h + images.h",
    ]
    items_joint = [
        "Confirm tensor memory layout with HW team: weight order [out_ch, in_ch, kH, kW]",
        "Agree on board input format: same Normalize(0.5, 0.5) as training",
        "On-board accuracy measurement vs 98.92 % SW baseline",
        "Runtime measurement: only preprocessing + kernel loop + accumulation",
    ]
    items_sw = [
        "Pynq board driver notebook (numpy / pynq only — course rule)",
        "Final accuracy vs runtime comparison table for report",
    ]

    def _section(title, col, items, left, top, width, height):
        _rect(s, left, top, width, Inches(0.38), col)
        _box(s, title, left+Inches(0.1), top, width, Inches(0.38),
             font_size=13, bold=True, color=C_DARK)
        y = top + Inches(0.42)
        for it in items:
            _box(s, "• " + it, left+Inches(0.1), y, width-Inches(0.2), Inches(0.35),
                 font_size=12, color=C_LIGHT)
            y += Inches(0.36)

    _section("Hardware (in progress)", C_RED,
             items_hw, Inches(0.35), Inches(1.55), Inches(4.0), Inches(2.3))
    _section("Joint agreement needed", C_YELLOW,
             items_joint, Inches(4.6), Inches(1.55), Inches(4.5), Inches(2.3))
    _section("SW remaining", C_GREEN,
             items_sw, Inches(9.3), Inches(1.55), Inches(3.7), Inches(1.2))

    _note(s,
        "On the hardware side the team is implementing the HLS kernel to match our forward graph order. "
        "The most important joint agreement is that the memory layout of each tensor and the input "
        "normalisation must exactly match. "
        "On the software side I still need to write the board-side Jupyter notebook driver "
        "using only NumPy and PYNQ as required by the course. "
        "After the bitstream is available we measure inference accuracy and runtime on the board "
        "and compare against our ninety-eight point ninety-two percent software baseline.")


def slide_qna(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _rect(s, 0, Inches(2.8), W, Inches(2.0), C_SURFACE)
    _box(s, "Thank You", Inches(0), Inches(2.9), W, Inches(0.9),
         font_size=52, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _box(s, "Questions?", Inches(0), Inches(3.8), W, Inches(0.6),
         font_size=28, color=C_WHITE, align=PP_ALIGN.CENTER)
    _box(s,
         "Checkpoint: QUANTIZED_MNISTCNN-2026-05-19-00-18-25.pth  |  "
         "Test acc: 98.92 %  |  Full-set: 99.36 %",
         Inches(0), Inches(5.4), W, Inches(0.4),
         font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)
    _note(s,
        "Thank you. I am happy to take questions on the quantization scheme, "
        "the exact forward-pass order, the export format, or the interface contract with the hardware team.")


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_cover(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_fixedpoint(prs)
    slide_qat_forward(prs)
    slide_training_setup(prs)
    slide_training_results(prs)
    slide_weight_embed(prs)
    slide_image_embed(prs)
    slide_pipeline(prs)
    slide_next(prs)
    slide_qna(prs)

    prs.save(OUT)
    print(f"Wrote  {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
